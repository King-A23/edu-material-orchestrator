from __future__ import annotations

from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pypdfium2 as pdfium
from pydantic import Field
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ...models import SerializableModel
from ...utils.files import ensure_directory
from ...utils.provenance import make_source_ref


class ExtractedImage(SerializableModel):
    source_ref: str
    image_path: str
    kind: str
    metadata: dict[str, str] = Field(default_factory=dict)


DOCX_NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
}


def extract_pptx_images(pptx_path: str | Path, output_dir: str | Path) -> list[ExtractedImage]:
    presentation = Presentation(str(pptx_path))
    assets_dir = ensure_directory(output_dir)
    extracted: list[ExtractedImage] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        source_ref = make_source_ref("pptx", slide_index)
        for image_index, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            extension = shape.image.ext or "png"
            image_path = assets_dir / f"slide_{slide_index:03d}_image_{image_index:03d}.{extension}"
            image_path.write_bytes(shape.image.blob)
            extracted.append(
                ExtractedImage(
                    source_ref=source_ref,
                    image_path=str(image_path),
                    kind="embedded_image",
                    metadata={"extension": extension},
                )
            )
    return extracted


def _docx_relationship_targets(archive: ZipFile) -> dict[str, str]:
    try:
        payload = archive.read("word/_rels/document.xml.rels")
    except KeyError:
        return {}

    root = ET.fromstring(payload)
    targets: dict[str, str] = {}
    for relationship in root.findall("rel:Relationship", DOCX_NAMESPACES):
        rel_type = relationship.attrib.get("Type", "")
        if not rel_type.endswith("/image"):
            continue
        rel_id = relationship.attrib.get("Id")
        target = relationship.attrib.get("Target")
        if rel_id and target:
            targets[rel_id] = target
    return targets


def _docx_media_target(target: str) -> str:
    normalized = target.replace("\\", "/")
    if normalized.startswith("/"):
        normalized = normalized[1:]
    if normalized.startswith("word/"):
        return normalized
    if normalized.startswith("../"):
        normalized = normalized[3:]
    return f"word/{normalized}"


def extract_docx_images(docx_path: str | Path, output_dir: str | Path) -> list[ExtractedImage]:
    assets_dir = ensure_directory(output_dir)
    extracted: list[ExtractedImage] = []
    image_counter = 0
    paragraph_index = 0

    with ZipFile(str(docx_path)) as archive:
        try:
            document_xml = archive.read("word/document.xml")
        except KeyError:
            return []
        rel_targets = _docx_relationship_targets(archive)
        root = ET.fromstring(document_xml)
        body = root.find("w:body", DOCX_NAMESPACES)
        if body is None:
            return []

        for child in body:
            local_name = child.tag.rsplit("}", maxsplit=1)[-1]
            if local_name == "p":
                paragraph_index += 1
                association = "inline_paragraph"
                target_index = paragraph_index
            elif local_name == "tbl":
                association = "nearest_paragraph"
                target_index = max(paragraph_index, 1)
            else:
                continue

            for blip in child.findall(".//a:blip", DOCX_NAMESPACES):
                rel_id = blip.get(f"{{{DOCX_NAMESPACES['r']}}}embed")
                if not rel_id:
                    continue
                target = rel_targets.get(rel_id)
                if not target:
                    continue
                media_target = _docx_media_target(target)
                try:
                    image_bytes = archive.read(media_target)
                except KeyError:
                    continue

                image_counter += 1
                extension = Path(media_target).suffix or ".bin"
                image_path = assets_dir / f"paragraph_{target_index:03d}_image_{image_counter:03d}{extension}"
                image_path.write_bytes(image_bytes)
                extracted.append(
                    ExtractedImage(
                        source_ref=make_source_ref("docx", target_index),
                        image_path=str(image_path),
                        kind="docx_image",
                        metadata={
                            "association": association,
                            "paragraph_index": str(target_index),
                            "relationship_id": rel_id,
                        },
                    )
                )
    return extracted


def render_pdf_pages_to_images(
    pdf_path: str | Path,
    output_dir: str | Path,
    prefix: str = "page",
) -> list[ExtractedImage]:
    assets_dir = ensure_directory(output_dir)
    pdf = pdfium.PdfDocument(str(pdf_path))
    extracted: list[ExtractedImage] = []
    try:
        for page_index in range(len(pdf)):
            page = pdf[page_index]
            bitmap = page.render(scale=2.0)
            pil_image = bitmap.to_pil()
            image_path = assets_dir / f"{prefix}_{page_index + 1:03d}.png"
            pil_image.save(image_path)
            extracted.append(
                ExtractedImage(
                    source_ref=make_source_ref("pdf", page_index + 1),
                    image_path=str(image_path),
                    kind="page_render",
                )
            )
            close_page = getattr(page, "close", None)
            if callable(close_page):
                close_page()
            close_bitmap = getattr(bitmap, "close", None)
            if callable(close_bitmap):
                close_bitmap()
    finally:
        close_pdf = getattr(pdf, "close", None)
        if callable(close_pdf):
            close_pdf()
    return extracted
