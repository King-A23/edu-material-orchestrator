from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from ...models.source import SourceUnit
from ...utils.hashing import make_manifest_id
from ...utils.provenance import make_source_ref
from ..common.base import DetectedInput, ReadError, ReaderOutput, SourceReader, build_source_document


def _clean_lines(lines: list[str]) -> str:
    return "\n".join(line.strip() for line in lines if line and line.strip())


def _extract_slide_text(slide) -> str:
    lines: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            lines.append(shape.text)
        elif getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if row_text:
                    lines.append(" | ".join(row_text))
    return _clean_lines(lines)


def _extract_notes_text(slide) -> str | None:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return None

    lines: list[str] = []
    for shape in notes_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip()
        if not text or text.lower() == "click to add notes":
            continue
        lines.append(text)
    notes_text = _clean_lines(lines)
    return notes_text or None


class PptxReader(SourceReader):
    supported_type = "pptx"

    def read(self, detected: DetectedInput) -> ReaderOutput:
        presentation_path = Path(detected.resolved_path)
        try:
            presentation = Presentation(presentation_path)
        except Exception as error:
            raise ReadError(f"Failed to open PPTX file '{presentation_path}': {error}") from error

        document = build_source_document(detected)
        units: list[SourceUnit] = []

        for index, slide in enumerate(presentation.slides, start=1):
            source_ref = make_source_ref("pptx", index)
            raw_text = _extract_slide_text(slide)
            notes_text = _extract_notes_text(slide)
            units.append(
                SourceUnit(
                    unit_id=make_manifest_id(document.id, source_ref),
                    source_id=document.id,
                    index=index,
                    kind="slide",
                    raw_text=raw_text,
                    notes_text=notes_text,
                    confidence=1.0,
                    source_ref=source_ref,
                )
            )

        return ReaderOutput(
            document=document,
            units=units,
            metadata={
                "reader": "open.pptx_reader",
                "slide_count": len(units),
                "notes_detected": sum(1 for unit in units if unit.notes_text),
            },
        )
