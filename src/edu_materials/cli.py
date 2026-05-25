from __future__ import annotations

import json
import shutil
import sys
from pathlib import Path
from typing import Annotated, Any, Literal

import typer

from .backends.open.ocr import detect_tesseract
from .backends.open.slide_renderer import detect_slide_renderer
from .backends.common.base import ReadError
from .config import AppConfig
from .pipeline.batch import build_batch, load_batch_inputs
from .pipeline.course_library import index_course_library_from_manifest, load_library_summary, query_course_library
from .pipeline.export_markdown import export_markdown_file
from .pipeline.ingest import inspect_source
from .pipeline.analyze_questions import AdapterInvocationError
from .pipeline.qa import run_qa_from_manifest
from .pipeline.qa_assignment import run_assignment_qa_from_manifest
from .pipeline.qa_quiz import run_quiz_qa_from_manifest
from .pipeline.render_docx import build_handout
from .pipeline.render_markdown import build_assignment_analysis
from .pipeline.render_quiz import build_quiz, discover_reference_inputs
from .pipeline.study_tools import (
    build_cram_pack,
    build_cram_plan,
    build_mistake_book,
    build_review_pack,
    build_rubric,
    build_variants,
    export_question_bank,
    grade_submission,
    refresh_mastery_artifacts,
)
from .utils.dependencies import build_capability_report
from .utils.installers import InstallOption, InstallationSummary, apply_installation_policy

try:
    from rich.console import Console
except ModuleNotFoundError:  # pragma: no cover - fallback for incomplete local envs
    Console = None


app = typer.Typer(
    add_completion=False,
    no_args_is_help=True,
    help="Build structured teaching handouts and assignment analyses from source material.",
)


def _console() -> Console | None:
    return Console() if Console is not None else None


def _emit_json(payload: dict[str, Any]) -> None:
    text = json.dumps(payload, ensure_ascii=False, indent=2)
    console = _console()
    if console is None:
        print(text)
        return
    console.print_json(text)


def _exit_with_usage_error(message: str) -> None:
    typer.echo(message, err=True)
    raise typer.Exit(code=2)


def _interactive_install_confirmation(issue, option: InstallOption) -> bool:
    command = " ".join(option.command)
    typer.echo(f"Missing dependency: {issue.display_name}")
    typer.echo(f"Reason: {issue.message}")
    typer.echo(f"Proposed command: {command}")
    if option.notes:
        typer.echo(option.notes)
    return typer.confirm("Install now?", default=True)


def _is_interactive_terminal() -> bool:
    return sys.stdin.isatty() and sys.stdout.isatty()


def _preferred_install_option(issue) -> InstallOption | None:
    for option in issue.install_options:
        if option.auto_supported and option.command:
            return option
    return None


def _noninteractive_dependency_guidance_lines(capability_report) -> list[str]:
    lines = [
        "Missing dependencies detected.",
        "Non-interactive session: --install-missing ask cannot prompt, so automatic installation was skipped.",
    ]

    for issue in capability_report.issues:
        requirement = "Required" if issue.is_required else "Optional"
        lines.append(f"{requirement} dependency: {issue.display_name}")
        lines.append(f"Reason: {issue.message}")

        preferred_option = _preferred_install_option(issue)
        preferred_notes = None
        if preferred_option is not None:
            lines.append(f"Suggested auto-install command: {' '.join(preferred_option.command)}")
            preferred_notes = preferred_option.notes
        else:
            lines.append("No automatic installer is available for this dependency in the current environment.")

        manual_notes = next(
            (option.notes for option in issue.install_options if option.notes and option.notes != preferred_notes),
            None,
        )
        if manual_notes:
            lines.append(f"Manual setup: {manual_notes}")

    if any(issue.is_required for issue in capability_report.issues):
        lines.append(
            "Required dependencies are still missing. Rerun with --install-missing auto or install them manually."
        )
    else:
        lines.append(
            "Continuing without the optional capabilities above. Rerun with --install-missing auto or install them manually."
        )
    return lines


def _emit_noninteractive_dependency_guidance(capability_report) -> None:
    for line in _noninteractive_dependency_guidance_lines(capability_report):
        typer.echo(line, err=True)


def _maybe_install_dependencies(
    capability_report,
    install_missing: str,
) -> InstallationSummary | None:
    if not capability_report.issues:
        return None

    effective_mode = install_missing
    if install_missing == "ask" and not _is_interactive_terminal():
        _emit_noninteractive_dependency_guidance(capability_report)
        effective_mode = "never"

    confirm_callback = _interactive_install_confirmation if effective_mode == "ask" else None
    return apply_installation_policy(
        capability_report.issues,
        mode=effective_mode,
        confirm_callback=confirm_callback,
        emit_callback=typer.echo,
    )


def _preflight_inputs(inputs: list[Path], target: str) -> tuple[list[Any], Any]:
    reader_outputs = [inspect_source(path) for path in inputs]
    capability_report = build_capability_report(target, reader_outputs)
    return reader_outputs, capability_report


def _installation_payload(summary: InstallationSummary | None) -> dict[str, Any] | None:
    return None if summary is None else summary.to_json_dict()


def _assignment_status(report) -> str:
    if (
        report.missing_source_ref_count
        or report.missing_answer_count
        or report.low_confidence_count
        or report.unclassified_count
        or report.missing_image_count
    ):
        return "warning"
    return "ok"


def _quiz_status(report) -> str:
    if (
        report.missing_answer_count
        or report.missing_explanation_count
        or report.missing_source_link_count
        or report.missing_image_count
    ):
        return "warning"
    return "ok"


def _load_config(config_path: str | Path | None = None) -> AppConfig:
    return AppConfig.load(config_path=config_path, cwd=Path.cwd())


def _quality_status(report, config: AppConfig) -> str:
    return "warning" if report.missing_source_ref_count > config.quality.max_missing_source_refs else "ok"


def _resolve_library_dir(config: AppConfig, course_dir: str | Path | None = None) -> Path:
    root = Path(course_dir).resolve() if course_dir is not None else Path.cwd().resolve()
    library_dir = root / config.library.root_dir_name
    if not library_dir.exists():
        raise FileNotFoundError(f"Course library does not exist: {library_dir}")
    return library_dir


def _write_demo_slide(placeholder, lines: list[str]) -> None:
    if not lines:
        placeholder.text = ""
        return
    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.text = lines[0]
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line


def _create_demo_pptx(output_dir: Path) -> Path:
    from pptx import Presentation

    demo_input = output_dir / "synthetic_agent_demo.pptx"
    presentation = Presentation()
    slides = [
        (
            "Agent-Ready Study Workflow",
            [
                "Input: PPTX, PDF, and DOCX course files",
                "Output: Markdown handout, DOCX handout, QA report",
                "Goal: keep source traceability visible",
            ],
        ),
        (
            "Why This Pipeline Stands Out",
            [
                "Local-first processing for ingestion and rendering",
                "Provider adapters stay optional",
                "Course-library artifacts can be reused later",
            ],
        ),
        (
            "Worked Example",
            [
                "Topic: derivatives and quick review loops",
                "Example: derivative of x^2 is 2x",
                "Checklist: inspect, build-handout, qa",
            ],
        ),
    ]

    for title, lines in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        _write_demo_slide(slide.placeholders[1], lines)

    presentation.save(demo_input)
    return demo_input


@app.command("inspect")
def inspect_command(
    input_path: Annotated[Path, typer.Argument(help="Source file to inspect.")],
    install_missing: Annotated[
        Literal["ask", "never", "auto"],
        typer.Option("--install-missing", help="Dependency installation policy."),
    ] = "never",
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Inspect a source file and output a manifest-like JSON payload."""

    _config = _load_config(config_path)
    try:
        result = inspect_source(input_path)
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    capability_report = build_capability_report("inspect", [result])
    installation_summary = _maybe_install_dependencies(capability_report, install_missing)
    if installation_summary is not None and installation_summary.executed_count > 0:
        result = inspect_source(input_path)
        capability_report = build_capability_report("inspect", [result])
    _emit_json(
        {
            **result.to_json_dict(),
            "capability_report": capability_report.to_json_dict(),
            "installation": _installation_payload(installation_summary),
        }
    )


@app.command("build-handout")
def build_handout_command(
    inputs: Annotated[
        list[Path],
        typer.Option(
            "--inputs",
            "-i",
            help="One or more source input paths.",
        ),
    ],
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    install_missing: Annotated[
        Literal["ask", "never", "auto"],
        typer.Option("--install-missing", help="Dependency installation policy."),
    ] = "ask",
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Optional explicit course directory for the local course library."),
    ] = None,
) -> None:
    """Build a Markdown-first handout and optionally derived export formats."""

    config = _load_config(config_path)
    if not inputs:
        raise typer.BadParameter("At least one --inputs value is required.")
    if output is None:
        output = config.paths.output_dir / "handout.md"
    try:
        _reader_outputs, capability_report = _preflight_inputs(inputs, "build-handout")
        installation_summary = _maybe_install_dependencies(capability_report, install_missing)
        if installation_summary is not None and installation_summary.executed_count > 0:
            _reader_outputs, capability_report = _preflight_inputs(inputs, "build-handout")
        if capability_report.missing_required_count > 0:
            _emit_json(
                {
                    "status": "error",
                    "reason": "missing_dependencies",
                    "capability_report": capability_report.to_json_dict(),
                    "installation": _installation_payload(installation_summary),
                }
            )
            raise typer.Exit(code=1)
        bundle, manifest, quality_report = build_handout(
            inputs,
            output,
            export_formats=export_to,
            config=config,
            course_dir=course_dir,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    status = _quality_status(quality_report, config)
    _emit_json(
        {
            "status": status,
            "output_bundle": bundle.to_json_dict(),
            "manifest": manifest.to_json_dict(),
            "quality_report": quality_report.to_json_dict(),
            "capability_report": capability_report.to_json_dict(),
            "installation": _installation_payload(installation_summary),
        }
    )


@app.command("demo")
def demo_command(
    output_dir: Annotated[
        Path | None,
        typer.Option("--output-dir", help="Directory where the synthetic input and demo outputs should be written."),
    ] = None,
    install_missing: Annotated[
        Literal["ask", "never", "auto"],
        typer.Option("--install-missing", help="Dependency installation policy."),
    ] = "never",
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Generate a redistributable synthetic PPTX and build a handout end to end."""

    config = _load_config(config_path)
    effective_output_dir = (output_dir or (config.paths.output_dir / "demo")).resolve()
    effective_output_dir.mkdir(parents=True, exist_ok=True)
    demo_input = _create_demo_pptx(effective_output_dir)
    demo_output = effective_output_dir / "demo_handout.md"

    try:
        _reader_outputs, capability_report = _preflight_inputs([demo_input], "build-handout")
        installation_summary = _maybe_install_dependencies(capability_report, install_missing)
        if installation_summary is not None and installation_summary.executed_count > 0:
            _reader_outputs, capability_report = _preflight_inputs([demo_input], "build-handout")
        if capability_report.missing_required_count > 0:
            _emit_json(
                {
                    "status": "error",
                    "reason": "missing_dependencies",
                    "demo_input": str(demo_input),
                    "capability_report": capability_report.to_json_dict(),
                    "installation": _installation_payload(installation_summary),
                }
            )
            raise typer.Exit(code=1)
        bundle, manifest, quality_report = build_handout(
            [demo_input],
            demo_output,
            export_formats=["docx"],
            config=config,
            course_dir=effective_output_dir,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error

    _emit_json(
        {
            "status": _quality_status(quality_report, config),
            "demo_input": str(demo_input),
            "output_bundle": bundle.to_json_dict(),
            "manifest": manifest.to_json_dict(),
            "quality_report": quality_report.to_json_dict(),
            "capability_report": capability_report.to_json_dict(),
            "installation": _installation_payload(installation_summary),
        }
    )


@app.command("qa")
def qa_command(
    manifest: Annotated[
        Path | None,
        typer.Option("--manifest", help="Manifest or intermediate JSON to validate."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Run QA from a build manifest and emit a quality report."""

    config = _load_config(config_path)
    if manifest is None:
        raise typer.BadParameter("--manifest is required.")
    try:
        report = run_qa_from_manifest(manifest)
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json(
        {
            "status": _quality_status(report, config),
            "quality_report": report.to_json_dict(),
        }
    )


@app.command("build-assignment-analysis")
def build_assignment_analysis_command(
    input_path: Annotated[
        Path | None,
        typer.Option("--input", help="Source assignment input path."),
    ] = None,
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    adapter_command: Annotated[
        str | None,
        typer.Option("--adapter-command", help="External adapter command used for question analysis."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    install_missing: Annotated[
        Literal["ask", "never", "auto"],
        typer.Option("--install-missing", help="Dependency installation policy."),
    ] = "ask",
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Optional explicit course directory for the local course library."),
    ] = None,
) -> None:
    """Build a Markdown-first assignment analysis and optionally derived export formats."""

    config = _load_config(config_path)
    if input_path is None:
        raise typer.BadParameter("--input is required.")
    if output is None:
        output = config.paths.output_dir / "assignment_analysis.md"
    effective_adapter_command = adapter_command or config.provider.adapter_command
    if not effective_adapter_command:
        typer.echo("--adapter-command is required.")
        raise typer.Exit(code=2)

    try:
        reader_outputs, capability_report = _preflight_inputs([input_path], "build-assignment-analysis")
        document_type = reader_outputs[0].document.type
        if document_type not in {"pdf", "docx", "pptx"}:
            raise typer.BadParameter(
                f"Unsupported assignment-analysis input type: {document_type}. Supported types are pdf, docx, and pptx."
            )
        installation_summary = _maybe_install_dependencies(capability_report, install_missing)
        if installation_summary is not None and installation_summary.executed_count > 0:
            reader_outputs, capability_report = _preflight_inputs([input_path], "build-assignment-analysis")
        bundle, manifest, quality_report = build_assignment_analysis(
            input_path,
            output,
            adapter_command=effective_adapter_command,
            export_formats=export_to,
            config=config,
            course_dir=course_dir,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError, AdapterInvocationError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error

    _emit_json(
        {
            "status": _assignment_status(quality_report),
            "output_bundle": bundle.to_json_dict(),
            "manifest": manifest.to_json_dict(),
            "quality_report": quality_report.to_json_dict(),
            "capability_report": capability_report.to_json_dict(),
            "installation": _installation_payload(installation_summary),
        }
    )


@app.command("qa-assignment")
def qa_assignment_command(
    manifest: Annotated[
        Path | None,
        typer.Option("--manifest", help="Assignment-analysis manifest to validate."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Run QA from an assignment-analysis manifest and emit a quality report."""

    _config = _load_config(config_path)
    if manifest is None:
        raise typer.BadParameter("--manifest is required.")
    try:
        report = run_assignment_qa_from_manifest(manifest)
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json(
        {
            "status": _assignment_status(report),
            "quality_report": report.to_json_dict(),
        }
    )


@app.command("build-quiz")
def build_quiz_command(
    references_dir: Annotated[
        Path | None,
        typer.Option("--references-dir", help="Directory of reference PDF/DOCX/PPTX files."),
    ] = None,
    manifest: Annotated[
        Path | None,
        typer.Option("--manifest", help="Optional YAML/JSON manifest of reference files."),
    ] = None,
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    prompt: Annotated[
        str | None,
        typer.Option("--prompt", help="Free-form quiz generation prompt."),
    ] = None,
    prompt_file: Annotated[
        Path | None,
        typer.Option("--prompt-file", help="Path to a file containing additional quiz prompt instructions."),
    ] = None,
    adapter_command: Annotated[
        str | None,
        typer.Option("--adapter-command", help="External adapter command used for quiz generation."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    install_missing: Annotated[
        Literal["ask", "never", "auto"],
        typer.Option("--install-missing", help="Dependency installation policy."),
    ] = "ask",
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Optional explicit course directory for the local course library."),
    ] = None,
) -> None:
    """Build a Markdown-first quiz from reference materials and a user prompt."""

    config = _load_config(config_path)
    if references_dir is None and manifest is None:
        raise typer.BadParameter("At least one of --references-dir or --manifest is required.")
    if not (prompt and prompt.strip()) and prompt_file is None:
        _exit_with_usage_error("At least one of --prompt or --prompt-file is required.")
    if output is None:
        output = config.paths.output_dir / "quiz.md"
    effective_adapter_command = adapter_command or config.provider.adapter_command
    if not effective_adapter_command:
        typer.echo("--adapter-command is required.")
        raise typer.Exit(code=2)

    try:
        reference_inputs = discover_reference_inputs(references_dir=references_dir, manifest_path=manifest)
        if not reference_inputs:
            raise typer.BadParameter("No supported reference inputs were found.")
        reference_paths = [Path(item["path"]).resolve() for item in reference_inputs]
        _reader_outputs, capability_report = _preflight_inputs(reference_paths, "build-quiz")
        installation_summary = _maybe_install_dependencies(capability_report, install_missing)
        if installation_summary is not None and installation_summary.executed_count > 0:
            _reader_outputs, capability_report = _preflight_inputs(reference_paths, "build-quiz")
        bundle, quiz_manifest, quality_report = build_quiz(
            output_path=output,
            adapter_command=effective_adapter_command,
            references_dir=references_dir,
            manifest_path=manifest,
            prompt=prompt,
            prompt_file=prompt_file,
            export_formats=export_to,
            config=config,
            course_dir=course_dir,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError, AdapterInvocationError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error

    _emit_json(
        {
            "status": _quiz_status(quality_report),
            "output_bundle": bundle.to_json_dict(),
            "manifest": quiz_manifest.to_json_dict(),
            "quality_report": quality_report.to_json_dict(),
            "capability_report": capability_report.to_json_dict(),
            "installation": _installation_payload(installation_summary),
        }
    )


@app.command("qa-quiz")
def qa_quiz_command(
    manifest: Annotated[
        Path | None,
        typer.Option("--manifest", help="Quiz manifest to validate."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Run QA from a quiz manifest and emit a quality report."""

    _config = _load_config(config_path)
    if manifest is None:
        raise typer.BadParameter("--manifest is required.")
    try:
        report = run_quiz_qa_from_manifest(manifest)
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json(
        {
            "status": _quiz_status(report),
            "quality_report": report.to_json_dict(),
        }
    )


@app.command("index-course-library")
def index_course_library_command(
    manifests: Annotated[
        list[Path],
        typer.Option("--manifest", help="One or more build manifests to index into the course library."),
    ] = [],
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Optional explicit course directory for the local course library."),
    ] = None,
) -> None:
    """Index one or more existing build manifests into a local course library."""

    config = _load_config(config_path)
    if not manifests:
        raise typer.BadParameter("At least one --manifest value is required.")
    try:
        library_dirs: list[str] = []
        for manifest in manifests:
            library_dir = index_course_library_from_manifest(manifest, config, course_dir=course_dir)
            library_dirs.append(str(library_dir))
        summary = load_library_summary(library_dirs[-1])
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json(
        {
            "status": "ok",
            "indexed_manifest_count": len(manifests),
            "library_dir": library_dirs[-1],
            "library_summary": summary,
        }
    )


@app.command("query-course-library")
def query_course_library_command(
    text: Annotated[
        str | None,
        typer.Option("--text", help="Optional free-text query across materials and questions."),
    ] = None,
    knowledge_point: Annotated[
        str | None,
        typer.Option("--knowledge-point", help="Optional canonical knowledge point filter."),
    ] = None,
    question_type: Annotated[
        str | None,
        typer.Option("--question-type", help="Optional question type filter."),
    ] = None,
    error_type: Annotated[
        str | None,
        typer.Option("--error-type", help="Optional error-type filter such as incorrect_answer or missing_key_steps."),
    ] = None,
    mastery_status: Annotated[
        str | None,
        typer.Option("--mastery-status", help="Optional mastery status filter: unseen, weak, developing, strong."),
    ] = None,
    review_only: Annotated[
        bool,
        typer.Option("--review-only/--no-review-only", help="Only return due-for-review mastery items."),
    ] = False,
    limit: Annotated[
        int,
        typer.Option("--limit", help="Maximum number of matched items to return."),
    ] = 10,
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Query the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        payload = query_course_library(
            library_dir,
            text=text,
            knowledge_point=knowledge_point,
            question_type=question_type,
            error_type=error_type,
            mastery_status=mastery_status,
            review_only=review_only,
            limit=limit,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json(payload)


@app.command("export-question-bank")
def export_question_bank_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    knowledge_point: Annotated[
        str | None,
        typer.Option("--knowledge-point", help="Optional knowledge point filter."),
    ] = None,
    question_type: Annotated[
        str | None,
        typer.Option("--question-type", help="Optional question type filter."),
    ] = None,
    limit: Annotated[
        int | None,
        typer.Option("--limit", help="Optional maximum number of questions to export."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Export a Markdown-first question bank from the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "question_bank.md"
        exported = export_question_bank(
            library_dir,
            effective_output,
            config=config,
            export_formats=export_to,
            knowledge_point=knowledge_point,
            question_type=question_type,
            limit=limit,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("build-mistake-book")
def build_mistake_book_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    attempts: Annotated[
        Path | None,
        typer.Option("--attempts", help="Optional JSON/YAML attempts file to import before building the mistake book."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build a Markdown-first mistake book from attempt records in the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "mistake_book.md"
        exported = build_mistake_book(
            library_dir,
            effective_output,
            config=config,
            attempts_path=attempts,
            export_formats=export_to,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("refresh-mastery")
def refresh_mastery_command(
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Refresh mastery, review queue, error taxonomy, and strategy patterns for the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        payload = refresh_mastery_artifacts(library_dir, config=config)
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", **payload})


@app.command("build-review-pack")
def build_review_pack_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    knowledge_point: Annotated[
        str | None,
        typer.Option("--knowledge-point", help="Optional knowledge point filter."),
    ] = None,
    limit: Annotated[
        int | None,
        typer.Option("--limit", help="Optional maximum number of review topics to include."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build a Markdown-first review pack from the current review queue."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "review_pack.md"
        exported = build_review_pack(
            library_dir,
            effective_output,
            config=config,
            export_formats=export_to,
            knowledge_point=knowledge_point,
            limit=limit,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("build-cram-pack")
def build_cram_pack_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    knowledge_point: Annotated[
        str | None,
        typer.Option("--knowledge-point", help="Optional knowledge point filter."),
    ] = None,
    limit: Annotated[
        int | None,
        typer.Option("--limit", help="Optional number of top knowledge points to include."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build a Markdown-first cram pack from the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "cram_pack.md"
        exported = build_cram_pack(
            library_dir,
            effective_output,
            config=config,
            export_formats=export_to,
            knowledge_point=knowledge_point,
            limit=limit,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("build-cram-plan")
def build_cram_plan_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    exam_date: Annotated[
        str | None,
        typer.Option("--exam-date", help="Exam date in YYYY-MM-DD format."),
    ] = None,
    days_available: Annotated[
        int | None,
        typer.Option("--days-available", help="Optional explicit number of study days."),
    ] = None,
    hours_per_day: Annotated[
        float | None,
        typer.Option("--hours-per-day", help="Optional hours available per day."),
    ] = None,
    knowledge_point: Annotated[
        str | None,
        typer.Option("--knowledge-point", help="Optional knowledge point filter."),
    ] = None,
    text: Annotated[
        str | None,
        typer.Option("--text", help="Optional free-text scope filter."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build a day-by-day cram plan from the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "cram_plan.md"
        exported = build_cram_plan(
            library_dir,
            effective_output,
            config=config,
            exam_date=exam_date,
            days_available=days_available,
            hours_per_day=hours_per_day,
            knowledge_point=knowledge_point,
            text=text,
            export_formats=export_to,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("build-rubric")
def build_rubric_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    question_id: Annotated[
        list[str],
        typer.Option("--question-id", help="Question IDs from the course library."),
    ] = [],
    question_record_id: Annotated[
        list[str],
        typer.Option("--question-record-id", help="Question record IDs from the course library."),
    ] = [],
    limit: Annotated[
        int | None,
        typer.Option("--limit", help="Optional maximum number of questions to include."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build a Markdown-first rubric from questions in the local course library."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "rubric.md"
        exported = build_rubric(
            library_dir,
            effective_output,
            config=config,
            question_ids=question_id,
            question_record_ids=question_record_id,
            export_formats=export_to,
            limit=limit,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("build-variants")
def build_variants_command(
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    question_id: Annotated[
        list[str],
        typer.Option("--question-id", help="Question IDs from the course library."),
    ] = [],
    question_record_id: Annotated[
        list[str],
        typer.Option("--question-record-id", help="Question record IDs from the course library."),
    ] = [],
    knowledge_point: Annotated[
        str | None,
        typer.Option("--knowledge-point", help="Optional knowledge point seed filter."),
    ] = None,
    count: Annotated[
        int | None,
        typer.Option("--count", help="Optional number of practice variants to generate."),
    ] = None,
    difficulty: Annotated[
        str | None,
        typer.Option("--difficulty", help="Optional difficulty label for adapter-based variants."),
    ] = None,
    adapter_command: Annotated[
        str | None,
        typer.Option("--adapter-command", help="Optional adapter command for synthetic variant generation."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build a Markdown-first variant or similar-practice pack from course-library questions."""

    config = _load_config(config_path)
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "variants.md"
        exported = build_variants(
            library_dir,
            effective_output,
            config=config,
            question_ids=question_id,
            question_record_ids=question_record_id,
            knowledge_point=knowledge_point,
            count=count,
            difficulty=difficulty,
            adapter_command=adapter_command or config.provider.adapter_command,
            export_formats=export_to,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError, AdapterInvocationError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("grade-submission")
def grade_submission_command(
    submission: Annotated[
        Path | None,
        typer.Option("--submission", help="Structured JSON/YAML submission file."),
    ] = None,
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Canonical Markdown output path, or a derived export path such as .docx/.pdf/.html."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Additional derived export formats: docx, pdf, html."),
    ] = [],
    adapter_command: Annotated[
        str | None,
        typer.Option("--adapter-command", help="Optional external adapter command used for model-based grading."),
    ] = None,
    course_dir: Annotated[
        Path | None,
        typer.Option("--course-dir", help="Course directory that owns the local course library."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Grade a structured submission against questions in the local course library."""

    config = _load_config(config_path)
    if submission is None:
        raise typer.BadParameter("--submission is required.")
    effective_adapter_command = adapter_command or config.provider.adapter_command
    try:
        library_dir = _resolve_library_dir(config, course_dir=course_dir)
        effective_output = output or config.paths.output_dir / "grading_report.md"
        exported, report = grade_submission(
            library_dir,
            submission,
            effective_output,
            config=config,
            export_formats=export_to,
            adapter_command=effective_adapter_command,
        )
    except (FileNotFoundError, IsADirectoryError, ValueError, ReadError, AdapterInvocationError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "grading_mode": report.grading_mode, "exports": exported, "grading_report": report.to_json_dict()})


@app.command("export")
def export_command(
    input_path: Annotated[
        Path | None,
        typer.Option("--input", help="Markdown input file."),
    ] = None,
    to: Annotated[
        list[str],
        typer.Option("--to", help="Export format(s): docx, pdf, html."),
    ] = [],
    output: Annotated[
        Path | None,
        typer.Option("--output", "-o", help="Optional explicit output path when exporting a single target."),
    ] = None,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Export a Markdown artifact to DOCX, PDF, or HTML."""

    config = _load_config(config_path)
    if input_path is None:
        raise typer.BadParameter("--input is required.")
    if not to:
        raise typer.BadParameter("At least one --to value is required.")
    try:
        exported = export_markdown_file(input_path, to, output_path=output, config=config)
    except (FileNotFoundError, ValueError) as error:
        typer.echo(str(error), err=True)
        raise typer.Exit(code=1) from error
    _emit_json({"status": "ok", "exports": exported})


@app.command("build-batch")
def build_batch_command(
    workflow: Annotated[
        Literal["handout", "assignment-analysis"],
        typer.Option("--workflow", help="Workflow to run for each batch input."),
    ] = "handout",
    inputs_dir: Annotated[
        Path | None,
        typer.Option("--inputs-dir", help="Directory of source inputs."),
    ] = None,
    manifest: Annotated[
        Path | None,
        typer.Option("--manifest", help="Optional YAML/JSON manifest of batch inputs."),
    ] = None,
    output_root: Annotated[
        Path | None,
        typer.Option("--output-root", help="Directory where per-input outputs should be written."),
    ] = None,
    adapter_command: Annotated[
        str | None,
        typer.Option("--adapter-command", help="External adapter command used for assignment-analysis."),
    ] = None,
    export_to: Annotated[
        list[str],
        typer.Option("--export-to", help="Derived export formats: docx, pdf, html."),
    ] = [],
    resume: Annotated[
        bool,
        typer.Option("--resume/--no-resume", help="Skip completed jobs whose manifest still matches the current input hash."),
    ] = True,
    config_path: Annotated[
        Path | None,
        typer.Option("--config", help="Optional config file path."),
    ] = None,
) -> None:
    """Build multiple handouts or assignment analyses from a directory or manifest."""

    config = _load_config(config_path)
    if inputs_dir is None and manifest is None:
        raise typer.BadParameter("At least one of --inputs-dir or --manifest is required.")
    items = load_batch_inputs(inputs_dir=inputs_dir, manifest_path=manifest)
    if not items:
        raise typer.BadParameter("No supported inputs were found for the requested batch.")
    effective_output_root = output_root or config.paths.output_dir / "batch"
    effective_adapter_command = adapter_command or config.provider.adapter_command
    if workflow == "assignment-analysis" and not effective_adapter_command:
        raise typer.BadParameter("--adapter-command is required for assignment-analysis batch builds.")
    summary = build_batch(
        workflow=workflow,
        items=items,
        output_root=effective_output_root,
        config=config,
        adapter_command=effective_adapter_command,
        export_formats=export_to,
        resume=resume,
    )
    _emit_json({"status": "ok", "summary": summary.to_json_dict()})


@app.command("doctor")
def doctor_command() -> None:
    """Inspect local runtime dependencies and supported adapter CLIs."""

    tesseract = detect_tesseract()
    slide_renderer = detect_slide_renderer()
    providers = {}
    for command_name in ("gemini", "claude", "codex"):
        providers[command_name] = {
            "available": shutil.which(command_name) is not None,
            "path": shutil.which(command_name),
        }

    _emit_json(
        {
            "status": "ok",
            "dependencies": {
                "tesseract": tesseract.to_json_dict(),
                "slide_renderer": slide_renderer.to_json_dict(),
            },
            "providers": providers,
        }
    )


def main() -> None:
    app()
