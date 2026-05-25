from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image, ImageDraw
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def create_sample_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Fractions Overview"
    slide.placeholders[1].text = "Definition\nEquivalent fractions\nWorked example 1/2 = 2/4"
    presentation.save(path)
    return path


def create_scanned_pdf(path: Path) -> Path:
    image = Image.new("RGB", (1400, 1800), "white")
    draw = ImageDraw.Draw(image)
    draw.text((80, 120), "Synthetic scanned worksheet", fill="black")
    draw.text((80, 220), "Fractions review: one half, one third, one quarter.", fill="black")
    draw.text((80, 320), "Example: 1/2 = 2/4", fill="black")
    image.save(path, "PDF", resolution=150.0)
    return path


def create_sample_image(path: Path, text: str = "Fixture image") -> Path:
    image = Image.new("RGB", (480, 280), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((20, 20, 460, 260), outline="black", width=3)
    draw.text((40, 120), text, fill="black")
    image.save(path, "PNG")
    return path


def create_assignment_pptx(path: Path, image_path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Homework 1"
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.0), Inches(3.5))
    textbox.text_frame.text = (
        "1. Calculate the derivative of x^2.\n"
        "(a) State the power rule.\n"
        "2. Explain Newton's first law in one sentence."
    )
    slide.shapes.add_picture(str(image_path), Inches(5.8), Inches(3.5), width=Inches(2.4))
    presentation.save(path)
    return path


def create_assignment_docx(path: Path, image_path: Path) -> Path:
    document = Document()
    document.add_heading("Assignment 1", level=1)
    document.add_paragraph("1. Explain the concept of inertia.")
    document.add_picture(str(image_path))
    document.add_paragraph("2. 设函数 f(x)=x+2，求 f(3)。")
    document.save(path)
    return path


def create_assignment_pdf(path: Path) -> Path:
    image = Image.new("RGB", (1600, 2000), "white")
    draw = ImageDraw.Draw(image)
    draw.text((80, 120), "1. What is 2 + 3?", fill="black")
    draw.text((80, 220), "2. Name one prime number greater than 10.", fill="black")
    draw.rectangle((1000, 1200, 1450, 1600), outline="black", width=4)
    draw.text((1060, 1380), "Graph", fill="black")
    image.save(path, "PDF", resolution=150.0)
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    return create_sample_pptx(tmp_path / "synthetic_fixture.pptx")


@pytest.fixture
def scanned_pdf(tmp_path: Path) -> Path:
    return create_scanned_pdf(tmp_path / "synthetic_scanned_fixture.pdf")


@pytest.fixture
def sample_image(tmp_path: Path) -> Path:
    return create_sample_image(tmp_path / "fixture_image.png")


@pytest.fixture
def assignment_pptx(tmp_path: Path, sample_image: Path) -> Path:
    return create_assignment_pptx(tmp_path / "assignment_fixture.pptx", sample_image)


@pytest.fixture
def assignment_docx(tmp_path: Path, sample_image: Path) -> Path:
    return create_assignment_docx(tmp_path / "assignment_fixture.docx", sample_image)


@pytest.fixture
def assignment_pdf(tmp_path: Path) -> Path:
    return create_assignment_pdf(tmp_path / "assignment_fixture.pdf")


@pytest.fixture
def mock_assignment_adapter_script() -> Path:
    return Path(__file__).resolve().parent / "helpers" / "mock_assignment_adapter.py"


@pytest.fixture
def quiz_references_dir(tmp_path: Path, sample_image: Path) -> Path:
    root = tmp_path / "quiz_refs"
    exam_dir = root / "exam"
    assignment_dir = root / "assignment"
    example_dir = root / "example"
    other_dir = root / "other"
    exam_dir.mkdir(parents=True)
    assignment_dir.mkdir(parents=True)
    example_dir.mkdir(parents=True)
    other_dir.mkdir(parents=True)

    create_assignment_docx(exam_dir / "exam_fixture.docx", sample_image)
    create_assignment_docx(assignment_dir / "assignment_fixture.docx", sample_image)
    create_assignment_pptx(example_dir / "example_fixture.pptx", sample_image)
    create_sample_pptx(other_dir / "notes_fixture.pptx")
    return root
